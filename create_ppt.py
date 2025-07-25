from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Create a presentation
prs = Presentation()
slide_width, slide_height = prs.slide_width, prs.slide_height

# Define title and bullet slide layout
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]

# Helper to add bullet slides
def add_bullet_slide(title, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = title
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.level = 0

# Slide 1: Project Overview
add_bullet_slide("Agentic RAG Chatbot for Multi-Format Document QA using MCP", [
    "Goal: Build a chatbot to answer questions using uploaded docs (PDF, DOCX, CSV, PPTX, TXT, MD).",
    "Agent-based architecture with message-passing (MCP).",
    "RAG: Ingest, Retrieve, Generate answer using context.",
    "Frontend built using Streamlit."
])

# Slide 2: Agentic Architecture with MCP
add_bullet_slide("Agent-Based Architecture with MCP", [
    "CoordinatorAgent: Orchestrates all agent interactions.",
    "IngestionAgent: Parses, chunks, embeds docs into vector store.",
    "RetrievalAgent: Retrieves top-k relevant chunks using FAISS.",
    "LLMResponseAgent: Uses retrieved context to generate final answer.",
    "Messages exchanged using structured JSON (MCP format)."
])

# Slide 3: System Flow Diagram (textual form)
add_bullet_slide("System Flow Diagram (Message Flow)", [
    "UI → CoordinatorAgent: USER_QUERY",
    "CoordinatorAgent → IngestionAgent: INGEST_REQUEST",
    "IngestionAgent → RetrievalAgent: RETRIEVE_REQUEST",
    "RetrievalAgent → LLMResponseAgent: CONTEXT_RESPONSE",
    "LLMResponseAgent → CoordinatorAgent → UI: ANSWER",
    "Trace_id used for tracking query lifecycle."
])

# Slide 4: Tech Stack
add_bullet_slide("Tech Stack Used", [
    "Frontend: Streamlit",
    "Embeddings: SentenceTransformers (all-MiniLM-L6-v2)",
    "Vector DB: FAISS",
    "LLM Backend: OpenAI GPT-4 / HuggingFace (mock fallback)",
    "Parsing: PyMuPDF, python-docx, pandas, pptx, markdown",
    "MCP Protocol: Custom in-memory JSON structure"
])

# Slide 5: UI Screenshots
slide = prs.slides.add_slide(blank_slide_layout)

# Add Title Textbox
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "📸 UI Screenshot"

# Add the image
image_path = r"D:\agentic_rag_mcp_bot\streamlit.png"  # Use raw string to avoid escape issues
left = Inches(1)
top = Inches(1.2)
height = Inches(4.5)  # Adjust as needed
slide.shapes.add_picture(image_path, left, top, height=height)


# Slide 6: Challenges & Future Scope
add_bullet_slide("Challenges Faced & Future Scope", [
    "Challenges: Parsing table-heavy documents, designing reusable MCP, chunking efficiency.",
    "Future Scope:",
    "- Add web/audio agents",
    "- Use Redis or Kafka for distributed MCP",
    "- Token streaming responses",
    "- Stateful multi-turn memory via trace_id"
])

# Save presentation

pptx_path = r"D:\agentic_rag_mcp_bot\presentation\Agentic_RAG_Chatbot_Presentation.pptx"
prs.save(pptx_path)
pptx_path
