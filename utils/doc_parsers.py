import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd
import markdown
import os


def parse_pdf(file_path):
    doc = fitz.open(file_path)
    text = "\n".join(page.get_text() for page in doc)
    return text

def parse_docx(file_path):
    doc = docx.Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def parse_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def parse_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def parse_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def parse_md(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        md_text = f.read()
    html = markdown.markdown(md_text)
    return html

def parse_document(file_path, file_type=None):
    ext = file_type or os.path.splitext(file_path)[-1].lower()
    if ext in ['.pdf']:
        text = parse_pdf(file_path)
    elif ext in ['.docx']:
        text = parse_docx(file_path)
    elif ext in ['.pptx']:
        text = parse_pptx(file_path)
    elif ext in ['.csv']:
        text = parse_csv(file_path)
    elif ext in ['.txt']:
        text = parse_txt(file_path)
    elif ext in ['.md', '.markdown']:
        text = parse_md(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    return {"text": text, "metadata": {"file_path": file_path, "file_type": ext}} 